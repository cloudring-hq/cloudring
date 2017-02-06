#!/usr/bin/env python3
"""Build an append-only, anonymized cloud-experience context from a mail archive.

The archive used by the first source is a ZIP-like container with local file
headers. Some tools cannot read its central directory, so this script streams
entries directly and never copies raw mail bodies into the knowledge folder.
"""

from __future__ import annotations

import argparse
import collections
import datetime as dt
import hashlib
import html
import io
import json
import math
import posixpath
import re
import struct
import sys
import textwrap
import xml.etree.ElementTree as ET
import zlib
from dataclasses import dataclass
from pathlib import Path
from typing import Callable, Iterable


try:
    import pypdf  # type: ignore
except Exception:  # pragma: no cover - optional runtime dependency
    pypdf = None

try:
    import docx  # type: ignore
except Exception:  # pragma: no cover
    docx = None

try:
    import openpyxl  # type: ignore
except Exception:  # pragma: no cover
    openpyxl = None

try:
    import pptx  # type: ignore
except Exception:  # pragma: no cover
    pptx = None


ZIP_LOCAL_FILE_HEADER = b"PK\x03\x04"


BRAND_OR_VENDOR_TERMS = [
    "1c",
    "admitad",
    "amazon",
    "amazon web services",
    "apple",
    "atlassian",
    "azure",
    "bitrix",
    "carrotquest",
    "cisco",
    "confluence",
    "docker",
    "facebook",
    "github",
    "gitlab",
    "google",
    "google cloud",
    "huawei",
    "hyper-v",
    "hostcms",
    "ibx",
    "infobox",
    "infoboxcloud",
    "instagram",
    "jira",
    "kubernetes",
    "linkedin",
    "mail.ru",
    "microsoft",
    "oceanstor",
    "office",
    "office365",
    "openstack",
    "oracle",
    "outlook",
    "plesk",
    "qiwi",
    "rzd",
    "sap",
    "sapsan",
    "selectel",
    "sharepoint",
    "skype",
    "slack",
    "telegram",
    "twitter",
    "ulmart",
    "viber",
    "visa",
    "vmware",
    "whatsapp",
    "yandex",
    "youtube",
    "яндекс",
    "битрикс",
    "инфобокс",
    "майкрософт",
    "сбер",
    "скайп",
    "телеграм",
]


DOMAIN_KEYWORDS: dict[str, list[str]] = {
    "compute_virtualization": [
        "blade",
        "cpu",
        "hypervisor",
        "instance",
        "ram",
        "server",
        "virtual",
        "vm",
        "vps",
        "vds",
        "виртуал",
        "гипервизор",
        "процессор",
        "сервер",
        "ядр",
    ],
    "storage_backup": [
        "backup",
        "block",
        "disk",
        "snapshot",
        "storage",
        "volume",
        "архив",
        "бэкап",
        "диск",
        "копия",
        "резерв",
        "снапшот",
        "хранилищ",
    ],
    "networking_connectivity": [
        "bgp",
        "dns",
        "ip",
        "load balanc",
        "nat",
        "network",
        "route",
        "ssl",
        "tcp",
        "vpn",
        "балансир",
        "днс",
        "доступ",
        "маршрут",
        "сеть",
        "сертификат",
    ],
    "identity_security": [
        "2fa",
        "account",
        "auth",
        "credential",
        "login",
        "password",
        "role",
        "security",
        "sso",
        "доступ",
        "логин",
        "парол",
        "пользовател",
        "роль",
        "секрет",
        "учет",
    ],
    "billing_commercial": [
        "act",
        "bill",
        "contract",
        "invoice",
        "payment",
        "price",
        "refund",
        "договор",
        "коммерческ",
        "оплат",
        "счет",
        "тариф",
        "цена",
    ],
    "support_incident": [
        "complaint",
        "escalat",
        "incident",
        "issue",
        "problem",
        "sla",
        "support",
        "ticket",
        "авар",
        "инцидент",
        "обращен",
        "поддерж",
        "проблем",
        "тикет",
    ],
    "developer_api": [
        "api",
        "automation",
        "cli",
        "git",
        "integration",
        "sdk",
        "webhook",
        "автоматизац",
        "интеграц",
        "скрипт",
    ],
    "platform_operations": [
        "deploy",
        "maintenance",
        "monitoring",
        "observability",
        "release",
        "upgrade",
        "доработ",
        "мониторинг",
        "обновлен",
        "релиз",
        "эксплуатац",
    ],
    "containers_clusters": [
        "cluster",
        "container",
        "namespace",
        "pod",
        "контейнер",
        "кластер",
        "оркестрац",
    ],
    "migration_onboarding": [
        "migration",
        "onboard",
        "transfer",
        "перенос",
        "миграц",
        "подключен",
        "переезд",
    ],
    "compliance_legal": [
        "audit",
        "compliance",
        "law",
        "license",
        "personal data",
        "policy",
        "аудит",
        "закон",
        "лиценз",
        "персональн",
        "политик",
        "регламент",
    ],
    "product_market": [
        "customer",
        "demo",
        "market",
        "offer",
        "pilot",
        "product",
        "sales",
        "клиент",
        "пилот",
        "продукт",
        "продаж",
        "рынок",
    ],
}


PROBLEM_KEYWORDS: dict[str, list[str]] = {
    "availability": [
        "down",
        "offline",
        "unavailable",
        "авар",
        "недоступ",
        "отказ",
        "паден",
    ],
    "latency_performance": [
        "delay",
        "latency",
        "slow",
        "timeout",
        "долго",
        "задерж",
        "медлен",
        "таймаут",
    ],
    "data_loss_recovery": [
        "backup",
        "corrupt",
        "lost",
        "restore",
        "восстанов",
        "потер",
        "резерв",
    ],
    "quota_capacity": [
        "capacity",
        "limit",
        "quota",
        "лимит",
        "квот",
        "мощност",
        "ресурс",
    ],
    "access_blocked": [
        "403",
        "access denied",
        "blocked",
        "forbidden",
        "permission",
        "блок",
        "доступ запрещ",
        "нет доступа",
    ],
    "billing_friction": [
        "act",
        "invoice",
        "overdue",
        "payment",
        "refund",
        "акт",
        "дебитор",
        "задолж",
        "оплат",
        "счет",
    ],
    "unclear_requirements": [
        "clarify",
        "requirements",
        "specification",
        "уточн",
        "требован",
        "формулиров",
    ],
    "integration_friction": [
        "api",
        "callback",
        "integration",
        "webhook",
        "интеграц",
        "обмен",
    ],
    "contractual_process": [
        "contract",
        "legal",
        "sign",
        "договор",
        "подпис",
        "юрид",
    ],
    "support_handoff": [
        "escalat",
        "forward",
        "support",
        "передал",
        "поддерж",
        "эскалац",
    ],
}


DOMAIN_IMPLICATIONS = {
    "compute_virtualization": [
        "Separate lifecycle state from capacity placement so failed starts, resizes, and migrations stay auditable.",
        "Expose resource requests, effective limits, and host-fit explanations before provisioning starts.",
    ],
    "storage_backup": [
        "Make backup, restore, snapshot retention, and failure evidence first-class tenant-facing workflows.",
        "Treat storage design as a product surface: capacity, latency, recovery points, and operational blast radius must be explicit.",
    ],
    "networking_connectivity": [
        "Offer self-service diagnostics for addresses, routes, names, certificates, and reachability.",
        "Keep tenant-facing network concepts stable even when provider-side topology changes.",
    ],
    "identity_security": [
        "Design identity flows around least privilege, recoverability, and visible audit trails.",
        "Avoid support-only access recovery paths for business-critical accounts.",
    ],
    "billing_commercial": [
        "Connect metering, invoices, legal documents, and service state so customers can reconcile usage without support.",
        "Commercial workflows need deterministic document generation and exception handling.",
    ],
    "support_incident": [
        "Incident handling must preserve timeline, ownership, customer-visible status, and next action.",
        "Support tooling should connect symptoms to platform telemetry instead of forcing manual correlation.",
    ],
    "developer_api": [
        "Every portal action should have an API path, idempotency, validation errors, and machine-readable state.",
        "Integrations should be observable and replayable when external systems fail.",
    ],
    "platform_operations": [
        "Provider operations need maintenance windows, rollbacks, version tracking, and impact communication.",
        "Internal operations should leave durable evidence that agents can reuse.",
    ],
    "containers_clusters": [
        "Cluster services require clear boundaries between tenant control, provider safety limits, and managed operations.",
        "Expose health and upgrade readiness rather than only raw orchestration objects.",
    ],
    "migration_onboarding": [
        "Onboarding should be modeled as a guided migration with inventory, prerequisites, tests, and rollback.",
        "Customers need predictable handoff from pre-sales promise to production operation.",
    ],
    "compliance_legal": [
        "Compliance artifacts should be generated from actual controls, policies, and audit events.",
        "Legal constraints must be visible in product limits and data handling choices.",
    ],
    "product_market": [
        "Packaging should connect customer job, operational cost, support burden, and measurable outcome.",
        "Pilot feedback must become reusable product requirements, not one-off promises.",
    ],
}


SCENARIO_BY_DOMAIN = {
    "compute_virtualization": "tenant capacity request or virtual compute lifecycle",
    "storage_backup": "persistent data protection or storage lifecycle",
    "networking_connectivity": "network reachability, naming, or secure connectivity",
    "identity_security": "access, account ownership, or security boundary",
    "billing_commercial": "commercial reconciliation, billing, or contract workflow",
    "support_incident": "support escalation or production incident workflow",
    "developer_api": "automation or integration workflow",
    "platform_operations": "provider-side operations and change management",
    "containers_clusters": "managed cluster or container platform workflow",
    "migration_onboarding": "customer onboarding or migration workflow",
    "compliance_legal": "policy, compliance, or legal process",
    "product_market": "product packaging, pilot, or market feedback",
}


ORG_LEGAL_RE = re.compile(
    r"\b(?:ООО|ОАО|АО|ЗАО|ПАО|ИП|LLC|Ltd|Inc|GmbH|S\.A\.|BV)\s+"
    r"(?:[\"'«][^\"'»]{1,80}[\"'»]|\S{2,80})",
    re.IGNORECASE,
)
EMAIL_RE = re.compile(r"[\w.+-]+@[\w.-]+\.[A-Za-zА-Яа-я]{2,}")
LOOSE_EMAIL_RE = re.compile(
    r"[\w.+-]+\s*@\s*[\w.-]+(?:\s*\.\s*[A-Za-zА-Яа-я]{2,})+",
    re.IGNORECASE,
)
URL_RE = re.compile(r"\b(?:https?|ftp)://[^\s<>\"]+|www\.[^\s<>\"]+", re.IGNORECASE)
DOMAIN_RE = re.compile(r"\b(?:[A-Za-z0-9-]{2,}\.)+[A-Za-zА-Яа-я]{2,}\b")
IP_RE = re.compile(r"\b(?:\d{1,3}\.){3}\d{1,3}\b")
PHONE_RE = re.compile(r"(?<!\w)(?:\+?\d[\d\s().-]{7,}\d)(?!\w)")
TICKET_RE = re.compile(r"\b(?:CRM|INC|REQ|SR|CASE|TASK|ID)[-_ #:]?\d{4,}\b", re.IGNORECASE)
LONG_ID_RE = re.compile(r"\b[A-Fa-f0-9]{16,}\b")
SECRET_LINE_RE = re.compile(
    r"(?im)^\s*(?:password|pass|passwd|пароль|login|логин|token|secret|ключ|key|ssh|authorized_keys)\b.*$"
)
PERSON_RU_RE = re.compile(r"\b[А-ЯЁ][а-яё]{2,}\s+[А-ЯЁ][а-яё]{2,}(?:\s+[А-ЯЁ][а-яё]{2,})?\b")
PERSON_EN_RE = re.compile(r"\b[A-Z][a-z]{2,}\s+[A-Z][a-z]{2,}(?:\s+[A-Z][a-z]{2,})?\b")


def utc_now_iso() -> str:
    return dt.datetime.now(dt.timezone.utc).replace(microsecond=0).isoformat()


def stable_hash(value: str | bytes, length: int = 12) -> str:
    if isinstance(value, str):
        value = value.encode("utf-8", errors="replace")
    return hashlib.sha256(value).hexdigest()[:length]


def parse_zip64_extra(extra: bytes, compressed_size: int, uncompressed_size: int) -> tuple[int, int]:
    offset = 0
    while offset + 4 <= len(extra):
        header_id, data_size = struct.unpack("<HH", extra[offset : offset + 4])
        offset += 4
        data = extra[offset : offset + data_size]
        offset += data_size
        if header_id != 0x0001:
            continue
        values = [
            struct.unpack("<Q", data[index : index + 8])[0]
            for index in range(0, (len(data) // 8) * 8, 8)
        ]
        if uncompressed_size == 0xFFFFFFFF and values:
            uncompressed_size = values.pop(0)
        if compressed_size == 0xFFFFFFFF and values:
            compressed_size = values.pop(0)
    return compressed_size, uncompressed_size


@dataclass(frozen=True)
class ArchiveEntry:
    index: int
    name: str
    method: int
    compressed_size: int
    uncompressed_size: int
    data_offset: int


def iter_archive_entries(
    archive_path: Path,
    should_read: Callable[[str, int], bool] | None = None,
) -> Iterable[tuple[ArchiveEntry, bytes | None]]:
    with archive_path.open("rb") as handle:
        index = 0
        while True:
            sig = handle.read(4)
            if not sig:
                return
            if sig != ZIP_LOCAL_FILE_HEADER:
                raise ValueError(f"Unexpected archive marker at entry {index}: {sig!r}")
            header = handle.read(26)
            if len(header) != 26:
                raise ValueError(f"Truncated local header at entry {index}")
            (
                _version_needed,
                flags,
                method,
                _mod_time,
                _mod_date,
                _crc32,
                compressed_size,
                uncompressed_size,
                name_len,
                extra_len,
            ) = struct.unpack("<HHHHHIIIHH", header)
            raw_name = handle.read(name_len)
            name = raw_name.decode("utf-8" if flags & 0x800 else "cp437", errors="replace")
            extra = handle.read(extra_len)
            compressed_size, uncompressed_size = parse_zip64_extra(
                extra,
                compressed_size,
                uncompressed_size,
            )
            data_offset = handle.tell()
            entry = ArchiveEntry(
                index=index,
                name=name,
                method=method,
                compressed_size=compressed_size,
                uncompressed_size=uncompressed_size,
                data_offset=data_offset,
            )
            data: bytes | None = None
            if should_read and should_read(name, compressed_size):
                compressed_data = handle.read(compressed_size)
                data = decompress_entry(method, compressed_data)
            else:
                handle.seek(compressed_size, io.SEEK_CUR)
            index += 1
            yield entry, data


def decompress_entry(method: int, data: bytes) -> bytes:
    if method == 0:
        return data
    if method == 8:
        return zlib.decompress(data, -15)
    raise ValueError(f"Unsupported archive compression method: {method}")


def strip_html(value: str) -> str:
    value = re.sub(r"(?is)<(script|style).*?>.*?</\1>", " ", value)
    value = re.sub(r"(?i)<br\s*/?>", "\n", value)
    value = re.sub(r"(?i)</(?:p|div|tr|li|h[1-6])>", "\n", value)
    value = re.sub(r"(?s)<[^>]+>", " ", value)
    value = html.unescape(value)
    return normalize_space(value)


def normalize_space(value: str) -> str:
    value = value.replace("\xa0", " ")
    value = re.sub(r"[ \t\r\f\v]+", " ", value)
    value = re.sub(r"\n\s+", "\n", value)
    value = re.sub(r"\n{3,}", "\n\n", value)
    return value.strip()


def build_brand_regex() -> re.Pattern[str]:
    escaped = [re.escape(term) for term in sorted(BRAND_OR_VENDOR_TERMS, key=len, reverse=True)]
    return re.compile(r"(?<![A-Za-z0-9])(?:" + "|".join(escaped) + r")(?![A-Za-z0-9])", re.IGNORECASE)


BRAND_RE = build_brand_regex()


def sanitize_text(value: str, max_chars: int | None = None) -> str:
    value = html.unescape(value or "")
    value = value.replace("©", "").replace("®", "").replace("™", "")
    replacements = [
        (SECRET_LINE_RE, "[secret-line]"),
        (URL_RE, "[url]"),
        (LOOSE_EMAIL_RE, "[email]"),
        (EMAIL_RE, "[email]"),
        (IP_RE, "[ip]"),
        (PHONE_RE, "[phone]"),
        (TICKET_RE, "[ticket]"),
        (ORG_LEGAL_RE, "[organization]"),
        (BRAND_RE, "[vendor]"),
        (DOMAIN_RE, "[domain]"),
        (LONG_ID_RE, "[id]"),
        (PERSON_RU_RE, "[person]"),
        (PERSON_EN_RE, "[person]"),
    ]
    for pattern, repl in replacements:
        value = pattern.sub(repl, value)
    value = re.sub(r"\[[a-z_]+]\s*(?:[,;/|]\s*)?\[[a-z_]+]", lambda m: m.group(0), value)
    value = re.sub(r"(?:\[email\][,; ]*){2,}", "[email-list] ", value)
    value = re.sub(r"(?:\[vendor\]\s*){2,}", "[vendor] ", value)
    value = re.sub(r"(?:\[domain\]\s*){2,}", "[domain] ", value)
    value = value.replace("@", "[at]")
    value = normalize_space(value)
    if max_chars and len(value) > max_chars:
        value = value[:max_chars].rstrip() + " ..."
    return value


def safe_subject_hint(value: str, max_chars: int = 160) -> str:
    value = sanitize_text(value, max_chars=max_chars)
    value = re.sub(r"\b[A-Z][A-Za-z0-9_-]{2,}\b", "[entity]", value)
    value = re.sub(r"\b[А-ЯЁ][а-яё]{2,}\s+[А-ЯЁ][а-яё]{1,}\b", "[person]", value)
    value = re.sub(
        r"\b(?:Юрий|Yuri|Анна|Владимир|Александр|Максим)\b",
        "[person]",
        value,
        flags=re.IGNORECASE,
    )
    return normalize_space(value)


def attachment_kind(extension: str, content_type: str) -> str:
    ext = extension.lower().strip(".")
    ctype = content_type.lower()
    if ext in {"png", "jpg", "jpeg", "gif", "bmp", "svg", "webp"} or ctype.startswith("image/"):
        return "image_or_diagram"
    if ext == "pdf" or "pdf" in ctype:
        return "document_pdf"
    if ext in {"doc", "docx", "odt"}:
        return "text_document"
    if ext in {"xls", "xlsx", "xlsm", "ods", "csv", "tsv"}:
        return "spreadsheet_or_table"
    if ext in {"ppt", "pptx", "odp"}:
        return "presentation"
    if ext in {"ics"}:
        return "calendar_invite"
    if ext in {"txt", "html", "htm", "xml", "json"} or ctype.startswith("text/"):
        return "text_or_markup"
    if ext in {"zip", "rar", "7z", "gz"}:
        return "archive_bundle"
    return "binary_or_unknown"


def text_of(root: ET.Element, tag: str) -> str:
    found = root.find(f".//{tag}")
    return found.text if found is not None and found.text else ""


def address_count(root: ET.Element, container_tag: str) -> int:
    found = root.find(f".//{container_tag}")
    if found is None:
        return 0
    return len(found.findall(".//emailAddress"))


def parse_boolish(value: str) -> bool:
    return value.strip().lower() in {"1", "1e0", "true", "yes"}


def parse_date(value: str) -> str:
    value = (value or "").strip()
    if not value:
        return ""
    return value[:19]


def period_from_date(value: str) -> str:
    if len(value) >= 7:
        return value[:7]
    return "unknown"


def body_text_from_message(root: ET.Element) -> str:
    body = text_of(root, "OPFMessageCopyHTMLBody") or text_of(root, "OPFMessageCopyBody")
    if "<" in body and ">" in body:
        body = strip_html(body)
    else:
        body = normalize_space(body)
    if not body:
        body = text_of(root, "OPFMessageCopyPreview")
    return body


def classify(text: str) -> tuple[list[str], list[str]]:
    folded = text.lower()
    domains = [
        name
        for name, words in DOMAIN_KEYWORDS.items()
        if any(word.lower() in folded for word in words)
    ]
    problems = [
        name
        for name, words in PROBLEM_KEYWORDS.items()
        if any(word.lower() in folded for word in words)
    ]
    if not domains:
        domains = ["general_business_context"]
    if not problems:
        problems = ["context_signal"]
    return domains, problems


def sentence_candidates(text: str, limit: int = 5) -> list[str]:
    pieces = re.split(r"(?<=[.!?])\s+|\n+", text)
    candidates: list[str] = []
    signal_words = [
        *sum(DOMAIN_KEYWORDS.values(), []),
        *sum(PROBLEM_KEYWORDS.values(), []),
    ]
    for piece in pieces:
        piece = normalize_space(piece)
        if len(piece) < 40:
            continue
        folded = piece.lower()
        score = sum(1 for word in signal_words if word.lower() in folded)
        if score:
            candidates.append(piece)
        if len(candidates) >= limit:
            break
    if not candidates and text:
        candidates = [text[:400]]
    return candidates[:limit]


def normalize_subject(subject: str) -> str:
    subject = re.sub(r"(?i)^\s*(re|fw|fwd|ответ|пересл):\s*", "", subject).strip()
    subject = TICKET_RE.sub("", subject)
    subject = LONG_ID_RE.sub("", subject)
    return normalize_space(subject.lower())


def mailbox_area(path: str) -> str:
    folded = path.lower()
    if "/sent items/" in folded:
        return "outgoing"
    if "/deleted items/" in folded:
        return "deleted"
    if "/inbox/" in folded:
        return "incoming"
    if "/calendar/" in folded:
        return "calendar"
    return "custom_folder"


def parse_attachment_size(value: str) -> int | None:
    value = (value or "").strip().replace(",", ".")
    if not value:
        return None
    try:
        parsed = float(value)
    except ValueError:
        return None
    if math.isnan(parsed) or math.isinf(parsed):
        return None
    return int(parsed)


def parse_message_xml(entry: ArchiveEntry, data: bytes) -> tuple[dict, list[dict]]:
    root = ET.fromstring(data)
    subject_raw = text_of(root, "OPFMessageCopySubject")
    thread_topic_raw = text_of(root, "OPFMessageCopyThreadTopic")
    body_raw = body_text_from_message(root)
    preview_raw = text_of(root, "OPFMessageCopyPreview") or body_raw[:500]
    sent_at = parse_date(text_of(root, "OPFMessageCopySentTime"))
    received_at = parse_date(text_of(root, "OPFMessageCopyReceivedTime"))
    date_value = sent_at or received_at or parse_date(text_of(root, "OPFMessageCopyModDate"))
    outgoing = parse_boolish(text_of(root, "OPFMessageIsOutgoing")) or mailbox_area(entry.name) == "outgoing"
    thread_raw = (
        text_of(root, "OPFMessageCopyLocalThreadGUID")
        or text_of(root, "OPFMessageCopyExchangeConversationId")
        or normalize_subject(subject_raw)
        or entry.name
    )
    subject_hint = safe_subject_hint(subject_raw)
    full_for_classification = f"{subject_raw}\n{preview_raw}\n{body_raw}"
    domains, problems = classify(full_for_classification)
    message_id = "msg-" + stable_hash(entry.name + text_of(root, "OPFMessageCopyMessageID"))
    attachments: list[dict] = []
    for index, attach in enumerate(root.findall(".//messageAttachment")):
        attrs = attach.attrib
        raw_url = attrs.get("OPFAttachmentURL", "")
        ext = attrs.get("OPFAttachmentContentExtension", "").lower().strip(".")
        content_type = attrs.get("OPFAttachmentContentType", "")
        raw_name = attrs.get("OPFAttachmentName", "")
        attach_id = "att-" + stable_hash(f"{entry.name}:{index}:{raw_url}:{raw_name}")
        attachments.append(
            {
                "attachment_id": attach_id,
                "message_id": message_id,
                "url_hash": stable_hash(raw_url),
                "raw_url": raw_url,
                "extension": sanitize_text(ext, max_chars=24).lower(),
                "kind": attachment_kind(ext, content_type),
                "declared_size_bytes": parse_attachment_size(
                    attrs.get("OPFAttachmentContentFileSize", "")
                ),
                "name_hash": stable_hash(raw_name),
            }
        )
    message = {
        "message_id": message_id,
        "entry_hash": stable_hash(entry.name),
        "period": period_from_date(date_value),
        "date": date_value,
        "direction": "outgoing" if outgoing else "incoming_or_archived",
        "mailbox_area": mailbox_area(entry.name),
        "thread_key": "thread-" + stable_hash(thread_raw),
        "subject_hint": subject_hint,
        "domain_tags": domains,
        "problem_tags": problems,
        "has_attachments": bool(attachments),
        "attachment_count": len(attachments),
        "participant_counts": {
            "from": address_count(root, "OPFMessageCopyFromAddresses"),
            "to": address_count(root, "OPFMessageCopyToAddresses"),
            "cc": address_count(root, "OPFMessageCopyCCAddresses"),
        },
        "content_length_chars": len(body_raw),
        "source_integrity": {
            "raw_subject_hash": stable_hash(subject_raw),
            "raw_body_hash": stable_hash(body_raw),
        },
    }
    return message, attachments


def extract_text_from_attachment(extension: str, content_type: str, data: bytes, limit: int) -> tuple[str, str]:
    ext = extension.lower().strip(".")
    ctype = content_type.lower()
    try:
        if ext == "pdf" or "pdf" in ctype:
            if pypdf is None:
                return "", "pdf_dependency_missing"
            reader = pypdf.PdfReader(io.BytesIO(data))
            parts: list[str] = []
            for page in reader.pages[:30]:
                parts.append(page.extract_text() or "")
                if sum(len(part) for part in parts) > limit:
                    break
            return normalize_space("\n".join(parts))[:limit], "text_extracted"
        if ext in {"docx"}:
            if docx is None:
                return "", "docx_dependency_missing"
            document = docx.Document(io.BytesIO(data))
            parts = [para.text for para in document.paragraphs]
            for table in document.tables[:10]:
                for row in table.rows[:40]:
                    parts.append(" | ".join(cell.text for cell in row.cells[:12]))
            return normalize_space("\n".join(parts))[:limit], "text_extracted"
        if ext in {"pptx"}:
            if pptx is None:
                return "", "pptx_dependency_missing"
            presentation = pptx.Presentation(io.BytesIO(data))
            parts = []
            for slide in presentation.slides[:40]:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        parts.append(shape.text)
            return normalize_space("\n".join(parts))[:limit], "text_extracted"
        if ext in {"xlsx", "xlsm"}:
            if openpyxl is None:
                return "", "xlsx_dependency_missing"
            workbook = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
            parts = []
            for sheet in workbook.worksheets[:8]:
                parts.append(f"sheet {sheet.title}")
                for row in sheet.iter_rows(max_row=80, max_col=20, values_only=True):
                    values = [str(cell) for cell in row if cell is not None]
                    if values:
                        parts.append(" | ".join(values))
                    if sum(len(part) for part in parts) > limit:
                        break
            return normalize_space("\n".join(parts))[:limit], "text_extracted"
        if ext in {"txt", "csv", "tsv", "ics", "html", "htm", "xml", "json"} or ctype.startswith("text/"):
            decoded = decode_bytes(data)
            if ext in {"html", "htm"} or "html" in ctype:
                decoded = strip_html(decoded)
            return normalize_space(decoded)[:limit], "text_extracted"
    except Exception as exc:  # pragma: no cover - depends on external files
        return "", f"extract_failed:{type(exc).__name__}"
    return "", "unsupported_binary_or_media"


def decode_bytes(data: bytes) -> str:
    for encoding in ("utf-8", "utf-16", "cp1251", "cp1252"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def write_jsonl(path: Path, rows: Iterable[dict]) -> None:
    with path.open("w", encoding="utf-8", newline="\n") as handle:
        for row in rows:
            handle.write(json.dumps(row, ensure_ascii=False, sort_keys=True) + "\n")


def read_jsonl(path: Path) -> list[dict]:
    rows: list[dict] = []
    with path.open("r", encoding="utf-8") as handle:
        for line in handle:
            line = line.strip()
            if line:
                rows.append(json.loads(line))
    return rows


def append_only_run_dir(root: Path, base_run_id: str) -> Path:
    root.mkdir(parents=True, exist_ok=True)
    candidate = root / base_run_id
    if not candidate.exists():
        return candidate
    suffix = 2
    while True:
        candidate = root / f"{base_run_id}-{suffix:02d}"
        if not candidate.exists():
            return candidate
        suffix += 1


def make_case_summary(case: dict) -> tuple[str, str, list[str]]:
    domains = case["domain_tags"]
    problems = case["problem_tags"]
    primary_domain = domains[0] if domains else "general_business_context"
    scenario = SCENARIO_BY_DOMAIN.get(primary_domain, "business or operational context")
    primary_problem = problems[0] if problems else "context_signal"
    situation = (
        f"A recurring source thread captured {scenario}; the dominant signal is "
        f"{primary_problem.replace('_', ' ')}."
    )
    need = (
        "Future cloud-platform agents should preserve the customer job, operational "
        "state, commercial state, and next action as separate but linked facts."
    )
    implications: list[str] = []
    for domain in domains[:5]:
        implications.extend(DOMAIN_IMPLICATIONS.get(domain, []))
    if not implications:
        implications.append(
            "Keep the source traceable and classify later when stronger domain evidence appears."
        )
    return situation, need, list(dict.fromkeys(implications))[:6]


def build_cases(messages: list[dict], attachments: list[dict]) -> list[dict]:
    attachments_by_message = collections.Counter(row["message_id"] for row in attachments)
    by_thread: dict[str, list[dict]] = collections.defaultdict(list)
    for message in messages:
        by_thread[message["thread_key"]].append(message)
    cases: list[dict] = []
    for thread_key, rows in sorted(by_thread.items()):
        rows = sorted(rows, key=lambda row: (row.get("date") or "", row["message_id"]))
        domain_counter = collections.Counter(
            tag for row in rows for tag in row.get("domain_tags", [])
        )
        problem_counter = collections.Counter(
            tag for row in rows for tag in row.get("problem_tags", [])
        )
        domain_tags = [tag for tag, _ in domain_counter.most_common()]
        problem_tags = [tag for tag, _ in problem_counter.most_common()]
        best = max(
            rows,
            key=lambda row: (
                len(row.get("domain_tags", [])) + len(row.get("problem_tags", [])),
                len(row.get("subject_hint", "")),
            ),
        )
        case_id = "case-" + stable_hash(thread_key)
        draft = {
            "case_id": case_id,
            "thread_key": thread_key,
            "period_start": rows[0].get("period", "unknown"),
            "period_end": rows[-1].get("period", "unknown"),
            "message_count": len(rows),
            "attachment_count": sum(attachments_by_message[row["message_id"]] for row in rows),
            "domain_tags": domain_tags,
            "problem_tags": problem_tags,
            "representative_subject_hint": best.get("subject_hint", ""),
            "evidence_message_ids": [row["message_id"] for row in rows[:20]],
            "evidence_message_id_count": len(rows),
        }
        situation, need, implications = make_case_summary(draft)
        draft.update(
            {
                "anonymized_situation": situation,
                "user_or_business_need": need,
                "platform_implications": implications,
            }
        )
        cases.append(draft)
    return cases


def count_values(rows: Iterable[dict], field: str) -> collections.Counter:
    counter: collections.Counter = collections.Counter()
    for row in rows:
        value = row.get(field)
        if isinstance(value, list):
            counter.update(value)
        elif value:
            counter[value] += 1
    return counter


def markdown_table(counter: collections.Counter, first_col: str, second_col: str = "count") -> str:
    lines = [f"| {first_col} | {second_col} |", "| --- | ---: |"]
    for key, count in counter.most_common():
        lines.append(f"| `{key}` | {count} |")
    return "\n".join(lines)


def write_docs(
    root: Path,
    run_dir: Path,
    manifest: dict,
    messages: list[dict],
    attachments: list[dict],
    cases: list[dict],
) -> None:
    domain_counts = count_values(messages, "domain_tags")
    problem_counts = count_values(messages, "problem_tags")
    period_counts = count_values(messages, "period")
    top_cases = sorted(
        cases,
        key=lambda row: (
            row.get("message_count", 0) + min(row.get("attachment_count", 0), 20),
            row.get("period_end", ""),
        ),
        reverse=True,
    )[:20]

    (root / "README.md").write_text(
        textwrap.dedent(
            f"""\
            # Knowledge Context

            This folder is append-only context for future cloud-platform agents. It stores
            anonymized experience records, derived cases, and coverage reports. Raw source
            messages, addresses, domains, company names, product names, and provider names
            are intentionally not copied here.

            Current source runs:

            - `{run_dir.name}`: legacy mail archive import, {manifest["parsed_messages"]} parsed messages,
              {manifest["case_count"]} derived cases, {manifest["attachments_referenced"]} referenced attachments.

            Usage rule for agents: prefer `cloud-experience-context.md` for principles,
            `imports/{run_dir.name}/cases.jsonl` for case-level retrieval, and
            `imports/{run_dir.name}/messages.jsonl` only when case evidence needs a
            message-level trace.
            """
        ),
        encoding="utf-8",
        newline="\n",
    )

    (root / "cloud-experience-context.md").write_text(
        textwrap.dedent(
            f"""\
            # Cloud Experience Context

            Source coverage: {manifest["parsed_messages"]} anonymized messages were assigned
            into {manifest["case_count"]} durable cases. The import keeps every parsed
            message reachable through a case id so later agents can append refinements
            without replacing earlier context.

            ## Agent Principles

            - Treat every customer request as a linked state machine: intent, resource state,
              commercial state, support state, and evidence trail.
            - Prefer self-service workflows with clear validation, idempotency, rollback,
              and audit output.
            - Make provider operations explainable to customers without exposing internal
              implementation names or vendor-specific dependencies.
            - Preserve long-tail support signals. Billing, documents, access recovery,
              migration, and incident communication are part of the cloud product, not
              side processes.
            - Keep context vendor-neutral. If a source mentions a named company, product,
              provider, person, address, domain, or ticket, use only the anonymized role
              and the reusable operational lesson.

            ## Recurrent Domains

            {markdown_table(domain_counts, "domain")}

            ## Recurrent Problem Signals

            {markdown_table(problem_counts, "signal")}

            ## Product Implications

            - Capacity must be explainable before provisioning, not only after failures.
            - Storage and backup need customer-visible recovery objectives, evidence, and
              tests.
            - Network, naming, certificate, and reachability diagnostics should be
              available without a support escalation.
            - Access management needs recoverable ownership, least privilege, and audit
              trails.
            - Billing and contract workflows should be connected to usage and service
              lifecycle so commercial friction does not become technical support load.
            - Incidents need timeline, ownership, impact, next action, and post-incident
              learning in the same durable record.
            - Migration and onboarding should be run as guided workflows with prerequisites,
              validation, fallback, and handoff to operations.

            ## Largest Derived Cases

            """
        )
        + "\n".join(
            f"- `{case['case_id']}` ({case['message_count']} messages, "
            f"{case['attachment_count']} attachments): {case['anonymized_situation']}"
            for case in top_cases
        )
        + "\n",
        encoding="utf-8",
        newline="\n",
    )

    (root / "agent-ingestion-guide.md").write_text(
        textwrap.dedent(
            f"""\
            # Agent Ingestion Guide

            This context is designed for retrieval and future model training datasets.
            Do not reconstruct source identities. Do not add raw quoted mail text. Append
            new source runs under `imports/`.

            ## Retrieval Order

            1. Read `cloud-experience-context.md` for reusable platform lessons.
            2. Retrieve `imports/{run_dir.name}/cases.jsonl` by `domain_tags` and
               `problem_tags`.
            3. Use `messages.jsonl` only for anonymized evidence snippets and source
               integrity hashes.
            4. Use `attachments.jsonl` to understand whether a thread included documents,
               diagrams, spreadsheets, meetings, or binary evidence.

            ## Record Guarantees

            - Every parsed message has a stable `message_id`.
            - Every parsed message belongs to exactly one `case_id`.
            - Raw addresses, domains, names, provider names, product names, and source
              paths are not stored.
            - Unsupported or binary attachments are still counted and linked by anonymized
              attachment ids.
            """
        ),
        encoding="utf-8",
        newline="\n",
    )

    (run_dir / "coverage.md").write_text(
        textwrap.dedent(
            f"""\
            # Source Run Coverage

            - Archive size bytes: {manifest["source_size_bytes"]}
            - Archive content hash: `{manifest["source_sha256"]}`
            - Archive entries scanned: {manifest["archive_entries_scanned"]}
            - Message XML entries found: {manifest["message_xml_entries"]}
            - Messages parsed: {manifest["parsed_messages"]}
            - Messages failed: {manifest["failed_messages"]}
            - Derived cases: {manifest["case_count"]}
            - Referenced attachments: {manifest["attachments_referenced"]}
            - Attachment archive entries matched: {manifest["attachment_entries_matched"]}
            - Attachment text extracted: {manifest["attachment_text_extracted"]}
            - Attachment extraction skipped or unsupported: {manifest["attachment_text_not_extracted"]}

            ## Messages By Period

            {markdown_table(period_counts, "period")}

            ## Messages By Domain

            {markdown_table(domain_counts, "domain")}

            ## Messages By Problem Signal

            {markdown_table(problem_counts, "signal")}
            """
        ),
        encoding="utf-8",
        newline="\n",
    )


def scan_forbidden(root: Path) -> list[dict]:
    issues: list[dict] = []
    scanners = {
        "known_vendor_or_brand_term": build_brand_regex(),
        "email": EMAIL_RE,
        "loose_email": LOOSE_EMAIL_RE,
        "url": URL_RE,
        "domain": DOMAIN_RE,
        "windows_path": re.compile(r"[A-Za-z]:\\"),
        "at_symbol": re.compile(r"@"),
    }
    for path in root.rglob("*"):
        if not path.is_file():
            continue
        if path.suffix.lower() not in {".md", ".json", ".jsonl"}:
            continue
        text = path.read_text(encoding="utf-8", errors="replace")
        for scanner_name, scanner in scanners.items():
            for match in scanner.finditer(text):
                token = match.group(0)
                if scanner_name == "domain" and token.startswith("mail-archive-"):
                    continue
                if scanner_name == "domain" and token.lower() in {
                    "attachments.jsonl",
                    "archive-batch-context.md",
                    "archive-backlog-cases.jsonl",
                    "archive-backlog-context.md",
                    "archive-backlog-members.jsonl",
                    "archive-backlog-signals.jsonl",
                    "archive-cases.jsonl",
                    "archive-members.jsonl",
                    "archive-residual-cases.jsonl",
                    "archive-residual-context.md",
                    "archive-residual-followup-cases.jsonl",
                    "archive-residual-followup-context.md",
                    "archive-residual-followup-controls.jsonl",
                    "archive-residual-followup-signals.jsonl",
                    "archive-residual-members.jsonl",
                    "archive-residual-signals.jsonl",
                    "archive-signals.jsonl",
                    "binary-batch-context.md",
                    "binary-cases.jsonl",
                    "binary-signals.jsonl",
                    "binary-textlike-cases.jsonl",
                    "binary-textlike-context.md",
                    "binary-textlike-followup-cases.jsonl",
                    "binary-textlike-followup-context.md",
                    "binary-textlike-followup-controls.jsonl",
                    "binary-textlike-followup-signals.jsonl",
                    "binary-textlike-signals.jsonl",
                    "cases.jsonl",
                    "database-index-cases.jsonl",
                    "database-index-context.md",
                    "database-index-residual-cases.jsonl",
                    "database-index-residual-context.md",
                    "database-index-residual-controls.jsonl",
                    "database-index-residual-signals.jsonl",
                    "database-index-signals.jsonl",
                    "directory-cases.jsonl",
                    "directory-context.md",
                    "document-batch-context.md",
                    "document-backlog-cases.jsonl",
                    "document-backlog-context.md",
                    "document-backlog-signals.jsonl",
                    "document-cases.jsonl",
                    "document-content-followup-cases.jsonl",
                    "document-content-followup-context.md",
                    "document-content-followup-controls.jsonl",
                    "document-content-followup-signals.jsonl",
                    "document-repair-followup-cases.jsonl",
                    "document-repair-followup-context.md",
                    "document-repair-followup-controls.jsonl",
                    "document-repair-followup-signals.jsonl",
                    "document-residual-cases.jsonl",
                    "document-residual-context.md",
                    "document-residual-controls.jsonl",
                    "document-residual-signals.jsonl",
                    "document-signals.jsonl",
                    "files.jsonl",
                    "text-batch-context.md",
                    "text-backlog-cases.jsonl",
                    "text-backlog-context.md",
                    "text-backlog-signals.jsonl",
                    "text-decoder-followup-cases.jsonl",
                    "text-decoder-followup-context.md",
                    "text-decoder-followup-controls.jsonl",
                    "text-decoder-followup-signals.jsonl",
                    "text-residual-cases.jsonl",
                    "text-residual-context.md",
                    "text-residual-followup-cases.jsonl",
                    "text-residual-followup-context.md",
                    "text-residual-followup-controls.jsonl",
                    "text-residual-followup-signals.jsonl",
                    "text-residual-signals.jsonl",
                    "text-cases.jsonl",
                    "text-signals.jsonl",
                    "messages.jsonl",
                    "manifest.json",
                    "media-batch-context.md",
                    "media-backlog-cases.jsonl",
                    "media-backlog-context.md",
                    "media-backlog-signals.jsonl",
                    "media-cases.jsonl",
                    "media-deep-followup-cases.jsonl",
                    "media-deep-followup-context.md",
                    "media-deep-followup-controls.jsonl",
                    "media-deep-followup-signals.jsonl",
                    "media-signals.jsonl",
                    "privacy-scan.json",
                    "coverage.md",
                    "context.md",
                    "coverage-backlog.md",
                    "coverage-backlog.json",
                    "credential-artifact-cases.jsonl",
                    "credential-artifact-context.md",
                    "credential-artifact-signals.jsonl",
                    "credential-lifecycle-cases.jsonl",
                    "credential-lifecycle-context.md",
                    "credential-lifecycle-controls.jsonl",
                    "credential-lifecycle-signals.jsonl",
                    "readme.md",
                    "cloud-experience-context.md",
                    "agent-ingestion-guide.md",
                    "failed-messages.jsonl",
                    "failed-files.jsonl",
                }:
                    continue
                issues.append(
                    {
                        "kind": scanner_name,
                        "path": str(path.relative_to(root)).replace("\\", "/"),
                        "term_hash": stable_hash(token.lower()),
                        "offset": match.start(),
                    }
                )
    return issues


def sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def run_import(args: argparse.Namespace) -> Path:
    archive_path = Path(args.archive).resolve()
    output_root = Path(args.output).resolve()
    imports_root = output_root / "imports"
    run_id = args.run_id or "mail-archive-" + dt.date.today().strftime("%Y%m%d")
    run_dir = append_only_run_dir(imports_root, run_id)
    run_dir.mkdir(parents=True)

    messages: list[dict] = []
    attachments: list[dict] = []
    attachment_by_raw_url: dict[str, dict] = {}
    failed_messages: list[dict] = []
    entry_count = 0
    message_xml_entries = 0

    def read_message_xml(name: str, size: int) -> bool:
        return "com.microsoft.__Messages/" in name and name.endswith(".xml") and size <= args.max_message_bytes

    for entry, data in iter_archive_entries(archive_path, read_message_xml):
        entry_count += 1
        if "com.microsoft.__Messages/" not in entry.name or not entry.name.endswith(".xml"):
            continue
        message_xml_entries += 1
        if data is None:
            continue
        try:
            message, message_attachments = parse_message_xml(entry, data)
        except Exception as exc:
            failed_messages.append(
                {
                    "entry_hash": stable_hash(entry.name),
                    "error": type(exc).__name__,
                }
            )
            continue
        messages.append(message)
        for attachment in message_attachments:
            raw_url = attachment.pop("raw_url")
            attachments.append(attachment)
            if raw_url:
                attachment_by_raw_url[raw_url] = attachment

    attachment_entries_matched = 0
    attachment_text_extracted = 0
    attachment_text_not_extracted = 0
    attachment_lookup = {row["attachment_id"]: row for row in attachments}

    def read_attachment(name: str, size: int) -> bool:
        if name not in attachment_by_raw_url:
            return False
        if size > args.max_attachment_bytes:
            return False
        row = attachment_by_raw_url[name]
        ext = row.get("extension", "")
        kind = row.get("kind", "")
        return bool(ext in SUPPORTED_ATTACHMENT_EXTENSIONS or kind in {"document_pdf", "text_or_markup"})

    for entry, data in iter_archive_entries(archive_path, read_attachment):
        if entry.name not in attachment_by_raw_url:
            continue
        row = attachment_lookup[attachment_by_raw_url[entry.name]["attachment_id"]]
        attachment_entries_matched += 1
        row["archive_entry_hash"] = stable_hash(entry.name)
        row["actual_size_bytes"] = entry.uncompressed_size
        row["source_integrity"] = {"content_sha256": ""}
        if data is None:
            row["text_status"] = (
                "too_large"
                if entry.uncompressed_size > args.max_attachment_bytes
                else "not_read"
            )
            attachment_text_not_extracted += 1
            continue
        row["source_integrity"]["content_sha256"] = hashlib.sha256(data).hexdigest()
        extracted, status = extract_text_from_attachment(
            row.get("extension", ""),
            row.get("kind", ""),
            data,
            args.attachment_text_limit,
        )
        row["text_status"] = status
        if extracted:
            domains, problems = classify(extracted)
            row["domain_tags"] = domains
            row["problem_tags"] = problems
            attachment_text_extracted += 1
        else:
            attachment_text_not_extracted += 1

    for row in attachments:
        row.setdefault("archive_entry_hash", "")
        row.setdefault("actual_size_bytes", None)
        row.setdefault("text_status", "not_found_in_archive_pass")
        row.setdefault("source_integrity", {"content_sha256": ""})

    cases = build_cases(messages, attachments)
    case_by_thread = {case["thread_key"]: case["case_id"] for case in cases}
    for message in messages:
        message["case_id"] = case_by_thread[message["thread_key"]]
    manifest = {
        "created_at": utc_now_iso(),
        "source_kind": "legacy_mail_archive",
        "source_size_bytes": archive_path.stat().st_size,
        "source_sha256": sha256_file(archive_path),
        "archive_entries_scanned": entry_count,
        "message_xml_entries": message_xml_entries,
        "parsed_messages": len(messages),
        "failed_messages": len(failed_messages),
        "case_count": len(cases),
        "attachments_referenced": len(attachments),
        "attachment_entries_matched": attachment_entries_matched,
        "attachment_text_extracted": attachment_text_extracted,
        "attachment_text_not_extracted": attachment_text_not_extracted,
        "append_only_run_id": run_dir.name,
        "sanitization": {
            "raw_messages_copied": False,
            "raw_source_paths_copied": False,
            "brand_or_vendor_names_masked": True,
            "addresses_domains_people_masked": True,
        },
    }

    write_jsonl(run_dir / "messages.jsonl", messages)
    write_jsonl(run_dir / "attachments.jsonl", attachments)
    write_jsonl(run_dir / "cases.jsonl", cases)
    write_jsonl(run_dir / "failed-messages.jsonl", failed_messages)
    (run_dir / "manifest.json").write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2, sort_keys=True) + "\n",
        encoding="utf-8",
        newline="\n",
    )
    write_docs(output_root, run_dir, manifest, messages, attachments, cases)
    issues = scan_forbidden(output_root)
    (run_dir / "privacy-scan.json").write_text(
        json.dumps(
            {
                "created_at": utc_now_iso(),
                "privacy_or_vendor_hits": issues,
                "hit_count": len(issues),
            },
            ensure_ascii=False,
            indent=2,
            sort_keys=True,
        )
        + "\n",
        encoding="utf-8",
        newline="\n",
    )
    if issues:
        raise RuntimeError(
            f"Generated context still contains {len(issues)} privacy/vendor term hits"
        )
    return run_dir


SUPPORTED_ATTACHMENT_EXTENSIONS = {
    "csv",
    "docx",
    "htm",
    "html",
    "ics",
    "json",
    "pdf",
    "pptx",
    "tsv",
    "txt",
    "xlsx",
    "xlsm",
    "xml",
}


def parse_args(argv: list[str]) -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--archive", required=True, help="Path to the legacy mail archive")
    parser.add_argument(
        "--output",
        default="knowledge-context",
        help="Append-only knowledge context output directory",
    )
    parser.add_argument("--run-id", default="", help="Optional append-only run id")
    parser.add_argument("--max-message-bytes", type=int, default=25 * 1024 * 1024)
    parser.add_argument("--max-attachment-bytes", type=int, default=15 * 1024 * 1024)
    parser.add_argument("--attachment-text-limit", type=int, default=250_000)
    parser.add_argument("--attachment-excerpt-chars", type=int, default=1800)
    return parser.parse_args(argv)


def main(argv: list[str]) -> int:
    try:
        run_dir = run_import(parse_args(argv))
    except Exception as exc:
        print(f"import failed: {exc}", file=sys.stderr)
        return 1
    print(str(run_dir))
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
